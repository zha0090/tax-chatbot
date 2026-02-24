from __future__ import annotations

import logging
import struct
from pathlib import Path

import olefile
from pptx import Presentation

from .csv_parser import ParsedChunk

logger = logging.getLogger(__name__)

TEXT_BYTES_ATOM = 0x0FA0
TEXT_CHARS_ATOM = 0x0FA8


def parse_ppt(file_path: str | Path) -> list[ParsedChunk]:
    """Parse a PowerPoint file (.ppt or .pptx) into per-slide text chunks."""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        slides = _parse_pptx(file_path)
    elif suffix == ".ppt":
        slides = _parse_legacy_ppt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {suffix}")

    return _slides_to_chunks(slides, file_path.name)


def _parse_pptx(file_path: Path) -> list[dict]:
    """Extract text from a .pptx file using python-pptx."""
    prs = Presentation(str(file_path))
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                texts.append(_extract_table_text(shape.table))
        if texts:
            slides.append({"slide_num": slide_num, "text": "\n".join(texts)})
    return slides


def _extract_table_text(table) -> str:
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    return "\n".join(rows)


def _parse_legacy_ppt(file_path: Path) -> list[dict]:
    """Extract text from a legacy .ppt file via OLE binary record parsing.

    The PPT binary format stores text in TextBytesAtom (0x0FA0) and
    TextCharsAtom (0x0FA8) records. Both are decoded as latin-1 since
    many legacy PPT files pack ASCII bytes into TextCharsAtom records.
    """
    try:
        ole = olefile.OleFileIO(str(file_path))
    except Exception:
        logger.exception("Failed to open OLE file: %s", file_path)
        return []

    try:
        data = ole.openstream("PowerPoint Document").read()
    except Exception:
        logger.exception("No PowerPoint Document stream in %s", file_path)
        return []
    finally:
        ole.close()

    raw_texts = _extract_ppt_text_records(data)
    return _group_texts_into_slides(raw_texts)


def _extract_ppt_text_records(data: bytes) -> list[str]:
    """Walk the PPT record chain and extract all text records."""
    texts: list[str] = []
    pos = 0
    while pos < len(data) - 8:
        rec_ver_inst = struct.unpack_from("<H", data, pos)[0]
        rec_ver = rec_ver_inst & 0x0F
        rec_type = struct.unpack_from("<H", data, pos + 2)[0]
        rec_len = struct.unpack_from("<I", data, pos + 4)[0]

        if rec_len > len(data) - pos - 8:
            break

        if rec_ver == 0x0F:
            pos += 8
            continue

        if rec_type in (TEXT_BYTES_ATOM, TEXT_CHARS_ATOM) and rec_len > 0:
            raw = data[pos + 8 : pos + 8 + rec_len]
            text = raw.decode("latin-1", errors="ignore").strip()
            if len(text) > 1 and not _is_template_text(text):
                texts.append(text)

        pos += 8 + rec_len

    return texts


def _is_template_text(text: str) -> bool:
    """Filter out PowerPoint template/placeholder strings."""
    template_patterns = [
        "Click to edit Master",
        "Second levelThird level",
        "* ",
    ]
    return any(p in text for p in template_patterns) or text == "*"


def _group_texts_into_slides(texts: list[str]) -> list[dict]:
    """Group extracted text blocks into logical slides.

    Since legacy PPT binary parsing doesn't give explicit slide boundaries,
    each meaningful text block becomes its own slide entry.
    """
    slides = []
    for idx, text in enumerate(texts):
        slides.append({"slide_num": idx + 1, "text": text})
    return slides


def _slides_to_chunks(slides: list[dict], source_file: str) -> list[ParsedChunk]:
    """Convert slide data into ParsedChunk objects."""
    chunks: list[ParsedChunk] = []
    for idx, slide in enumerate(slides):
        chunks.append(
            ParsedChunk(
                text=slide["text"],
                metadata={"slide_num": slide["slide_num"]},
                source_file=source_file,
                source_type="ppt",
                chunk_index=idx,
            )
        )
    return chunks
